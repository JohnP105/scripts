import sys
import os
from pdfminer.high_level import extract_text
from pptx import Presentation


def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    try:
        text = extract_text(pdf_path)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""


def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint file."""
    try:
        presentation = Presentation(pptx_path)
        extracted_text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_text.append(shape.text.strip())

        return "\n\n".join(extracted_text)
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""


def extract_text_from_file(input_path, output_txt_path):
    """Determine file type and extract text accordingly."""
    if not os.path.exists(input_path):
        print(f"Error: File '{input_path}' not found.")
        return

    file_extension = os.path.splitext(input_path)[1].lower()

    if file_extension == ".pdf":
        extracted_text = extract_text_from_pdf(input_path)
    elif file_extension == ".pptx":
        extracted_text = extract_text_from_pptx(input_path)
    else:
        print(
            "Unsupported file format! Please provide a PDF or PowerPoint (.pptx) file."
        )
        return

    # Save extracted text
    with open(output_txt_path, "w", encoding="utf-8") as text_file:
        text_file.write(extracted_text)

    print(f"Text extracted and saved to {output_txt_path}")


if __name__ == "__main__":
    input_path = "ENGR_Diversity.pptx"
    output_txt_path = "output.txt"
    extract_text_from_file(input_path, output_txt_path)
